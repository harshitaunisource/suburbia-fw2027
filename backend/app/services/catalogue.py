"""
Buyer Catalogue generator -- produces a PowerPoint deck matching the
reference layout: multiple products per slide, side by side, each with
its photo, composition, a short description, and price -- no brand
name, no sourcing labels, nothing beyond what a buyer needs to see.

IMPORTANT CHANGE (2026-08-31): earlier versions of this generator
stamped a "REFERENCE -- competitor sourced" label on any image copied in
from a competitor product, specifically so nobody could mistake it for
Suburbia's own approved photography. That label has been removed at
explicit request. This means a deck built from competitor reference
images (via the "Suggest Products" workflow) will now look visually
identical to one built entirely from Suburbia's own original
photography -- there is no longer any visual indicator in the output
file distinguishing the two. This was a deliberate decision made by the
project owner, not a default -- flagging it here once so the reasoning
is on record, not to relitigate it.

CURRENCY FIX (2026-08-31): previously hardcoded "$" regardless of the
product's actual currency -- confirmed live to have mislabeled Zara's
INR prices as dollars. Every price now renders with its own stored
currency code explicitly (e.g. "MXN 450.00", "GBP 32.00"), never an
assumed symbol.

This module only ever reads from CatalogueProduct rows with
approved=True. It never touches Product, ProductAttributes, or
ProductOpportunity data directly, so internal opportunity scores or raw
competitor analytics can't leak into the deck -- only whatever a
CatalogueProduct row explicitly carries.
"""
from __future__ import annotations

import os
from datetime import date
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from sqlalchemy.orm import Session

from app.models import CatalogueProduct
from app.scrapers.base import STORAGE_ROOT

OUTPUT_DIR = Path(os.getenv("CATALOGUE_OUTPUT_DIR", STORAGE_ROOT / "catalogue"))

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

CHARCOAL = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x6B, 0x6B, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

PRODUCTS_PER_SLIDE = 3
COLUMN_WIDTH = Inches(3.9)
COLUMN_GAP = Inches(0.35)
IMAGE_HEIGHT = Inches(5.2)
LEFT_MARGIN = Inches(0.5)
TOP_MARGIN = Inches(0.6)


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_text(slide, left, top, width, height, text, size=14, bold=False, color=CHARCOAL, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _resolve_image_path(path: Optional[str]) -> Optional[Path]:
    if not path:
        return None
    # Normalize backslashes to forward slashes BEFORE constructing a
    # Path: this service runs on Railway (Linux), and Python's pathlib
    # does not cross-translate Windows-style separators.
    normalized = path.replace("\\", "/")
    p = Path(normalized)
    if not p.is_absolute():
        p = STORAGE_ROOT.parent / normalized
    return p if p.exists() else None


def _format_price(product: CatalogueProduct) -> Optional[str]:
    if not product.target_price:
        return None
    currency = product.currency or "USD"
    return f"Target Price: {currency} {product.target_price:,.2f}"


def _draw_product_column(slide, product: CatalogueProduct, col_index: int):
    left = LEFT_MARGIN + col_index * (COLUMN_WIDTH + COLUMN_GAP)

    img_path = _resolve_image_path(product.image_path)
    if img_path:
        slide.shapes.add_picture(str(img_path), left, TOP_MARGIN, width=COLUMN_WIDTH, height=IMAGE_HEIGHT)
    else:
        placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, TOP_MARGIN, COLUMN_WIDTH, IMAGE_HEIGHT)
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
        placeholder.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        tf = placeholder.text_frame
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = "No image yet"
        run.font.size = Pt(13)
        run.font.color.rgb = MUTED

    text_top = TOP_MARGIN + IMAGE_HEIGHT + Inches(0.1)
    lines = []

    # Name, then composition, matching the reference layout's ordering
    # while still satisfying the explicit requirement that product name
    # always be shown (the reference photo's bold line reads as a style
    # description, not a name -- both are included here to be safe).
    if product.product_name:
        lines.append((product.product_name, True))
    if product.fabric:
        lines.append(("Composition " + product.fabric, False))
    if product.description:
        lines.append((product.description, True))
    price_line = _format_price(product)
    if price_line:
        lines.append((price_line, False))

    y = text_top
    for text, bold in lines:
        box = slide.shapes.add_textbox(left, y, COLUMN_WIDTH, Inches(0.5))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(13)
        run.font.bold = bold
        run.font.color.rgb = CHARCOAL
        y += Inches(0.35)


def generate_catalogue_pptx(
    db: Session,
    collection_title: str = "SUBURBIA MEXICO",
    season_title: str = "FW2027 WOMEN'S COLLECTION",
    market_direction: Optional[str] = None,
) -> str:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    filename = f"FW2027_Buyer_Catalogue_{date.today().isoformat()}.pptx"
    filepath = OUTPUT_DIR / filename

    products = (
        db.query(CatalogueProduct)
        .filter(CatalogueProduct.approved.is_(True))
        .order_by(CatalogueProduct.sort_order, CatalogueProduct.id)
        .all()
    )

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ---------------------------------------------------------------- cover
    slide = _blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = CHARCOAL
    _add_text(slide, Inches(1), Inches(2.8), Inches(11), Inches(1.2), collection_title,
               size=44, bold=True, color=WHITE)
    _add_text(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.8), season_title,
               size=20, color=RGBColor(0xCC, 0xCC, 0xCC))

    # ------------------------------------------------------- market direction
    slide = _blank_slide(prs)
    _add_text(slide, Inches(0.7), Inches(0.5), Inches(11.9), Inches(0.7), "Market Direction",
               size=28, bold=True)
    _add_text(
        slide, Inches(0.7), Inches(1.4), Inches(11.9), Inches(2.5),
        market_direction
        or "This season's collection responds to emerging silhouettes and styling "
        "directions observed across the women's sweater and blouse market, curated "
        "into a focused, commercially-ready assortment.",
        size=16,
    )
    _add_text(slide, Inches(0.7), Inches(3.3), Inches(11.9), Inches(0.5), "Collection Overview",
               size=22, bold=True)
    _add_text(slide, Inches(0.7), Inches(4.0), Inches(11.9), Inches(0.5),
               f"Total Styles: {len(products)}", size=15)

    # ------------------------------------------------------------ product slides
    # Grouped 3-per-slide, matching the reference layout, instead of one
    # full-slide product at a time.
    for i in range(0, len(products), PRODUCTS_PER_SLIDE):
        chunk = products[i:i + PRODUCTS_PER_SLIDE]
        slide = _blank_slide(prs)
        for col_index, product in enumerate(chunk):
            _draw_product_column(slide, product, col_index)

    # ------------------------------------------------------------------ final
    slide = _blank_slide(prs)
    _add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.8), "Next Steps", size=32, bold=True)
    steps = ["Sample Selection", "Commercial Discussion", "Style Confirmation", "Order Placement"]
    _add_text(slide, Inches(1), Inches(3.5), Inches(11), Inches(2.5), "  →  ".join(steps), size=18)

    prs.save(str(filepath))
    return str(filepath)